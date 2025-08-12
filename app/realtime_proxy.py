#!/usr/bin/env python3
"""
FastAPI proxy to OpenAI Realtime WebRTC SDP exchange.
Keeps API key on the server and returns the SDP answer to the browser.

Endpoint:
  POST /sdp    Content-Type: application/sdp, Body: offer SDP
  -> forwards to OpenAI Realtime and returns answer SDP (text/plain)

Usage (dev):
  uvicorn app.realtime_proxy:app --host 127.0.0.1 --port 8787

For HTTPS (self-signed):
  uvicorn app.realtime_proxy:app --host 127.0.0.1 --port 8787 \
    --ssl-keyfile .cert/key.pem --ssl-certfile .cert/cert.pem
"""
from fastapi import FastAPI, Response, Request
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import base64
import io
from fastapi.middleware.cors import CORSMiddleware
import os
import requests
from typing import Any
import re
import logging
from dotenv import load_dotenv, find_dotenv
from datetime import datetime


# Load environment from nearest .env (project root) before reading keys
try:
    load_dotenv(find_dotenv())
except Exception:
    pass

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_REALTIME_MODEL = os.getenv("OPENAI_REALTIME_MODEL", "gpt-4o-realtime-preview-2024-12-17")
OPENAI_REALTIME_URL = f"https://api.openai.com/v1/realtime?model={OPENAI_REALTIME_MODEL}"
TAVILY_API_KEY = os.getenv("TAVILY_API_KEY")

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
log = logging.getLogger("realtime_proxy")

app = FastAPI(title="AlitaOS Realtime Proxy")

# Allow any localhost origin (any port, http/https)
# In dev, allow any origin; we do not use credentials in these calls.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/health")
async def health():
    return {"status": "ok", "model": OPENAI_REALTIME_MODEL}

@app.post("/sdp")
async def sdp(request: Request):
    if not OPENAI_API_KEY:
        return Response(content="OPENAI_API_KEY not configured", media_type="text/plain", status_code=500)

    offer_sdp = await request.body()
    try:
        log.info("/sdp: forwarding offer to OpenAI Realtime: %s", OPENAI_REALTIME_MODEL)
        resp = requests.post(
            OPENAI_REALTIME_URL,
            data=offer_sdp,
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/sdp",
                "OpenAI-Beta": "realtime=v1",
            },
            timeout=30,
        )
        log.info("/sdp: OpenAI response status=%s", resp.status_code)
        return Response(content=resp.text, media_type="application/sdp", status_code=resp.status_code)
    except Exception as e:
        log.exception("/sdp: proxy error: %s", e)
        return Response(content=f"Proxy error: {e}", media_type="text/plain", status_code=500)


# Explicit preflight handlers (defensive; CORSMiddleware should handle this)
@app.options("/sdp")
async def sdp_options():
    return Response(status_code=200)

@app.options("/tool")
async def tool_options():
    return Response(status_code=200)


class ToolPayload(BaseModel):
    name: str
    args: dict | None = None


# --- Time resolver helpers -------------------------------------------------
_CITY_TZ_MAP: dict[str, str] = {
    # US
    "new york": "America/New_York",
    "nyc": "America/New_York",
    "los angeles": "America/Los_Angeles",
    "la": "America/Los_Angeles",
    "san francisco": "America/Los_Angeles",
    "seattle": "America/Los_Angeles",
    "chicago": "America/Chicago",
    "austin": "America/Chicago",
    "denver": "America/Denver",
    # Europe
    "london": "Europe/London",
    "paris": "Europe/Paris",
    "berlin": "Europe/Berlin",
    "madrid": "Europe/Madrid",
    "rome": "Europe/Rome",
    # Asia
    "karachi": "Asia/Karachi",
    "pakistan": "Asia/Karachi",
    "delhi": "Asia/Kolkata",
    "mumbai": "Asia/Kolkata",
    "bangalore": "Asia/Kolkata",
    "kolkata": "Asia/Kolkata",
    "tokyo": "Asia/Tokyo",
    "seoul": "Asia/Seoul",
    "singapore": "Asia/Singapore",
    "hong kong": "Asia/Hong_Kong",
    "shanghai": "Asia/Shanghai",
    "beijing": "Asia/Shanghai",
    # Oceania
    "sydney": "Australia/Sydney",
    "auckland": "Pacific/Auckland",
    # Middle East
    "dubai": "Asia/Dubai",
    "riyadh": "Asia/Riyadh",
    # Misc
    "utc": "Etc/UTC",
    "gmt": "Etc/UTC",
}

# Be more permissive and allow trailing qualifiers like "right now" or punctuation
_TIME_REGEX = re.compile(
    r"\b(?:what\s+time\s+is\s+it\s+in|(?:(?:current|local)\s+)?time\s+in)\s+([A-Za-z ,\-]+?)(?:\b|\s+right\s+now\b|\s+today\b|[\.!?,]|$)",
    re.IGNORECASE,
)


def _extract_place_from_query(q: str) -> str | None:
    """Return place/city string if query looks like a time question."""
    if not q:
        return None
    m = _TIME_REGEX.search(q.strip())
    if m:
        return m.group(1).strip(" -,.")
    return None


def _map_place_to_tz(place: str) -> tuple[str, str] | None:
    """Best-effort mapping from place -> (canonical_place, IANA tz)."""
    if not place:
        return None
    key = place.lower().strip()
    # direct
    if key in _CITY_TZ_MAP:
        return (place, _CITY_TZ_MAP[key])
    # try substring match on known keys
    for k, tz in _CITY_TZ_MAP.items():
        if k in key:
            return (place, tz)
    return None


@app.post("/tool")
async def tool_exec(payload: ToolPayload):
    """Execute assistant tools server-side and return JSON results.

    Supported tools:
      - image.generate: {"name":"image.generate", "args":{"prompt":"..."}}
    """
    if not OPENAI_API_KEY:
        return JSONResponse({"success": False, "error": "OPENAI_API_KEY not configured"}, status_code=500)

    name = payload.name
    args = payload.args or {}
    try:
        log.info("/tool: name=%s args=%s", name, args)
    except Exception:
        pass

    if name == "image.generate":
        prompt = args.get("prompt", "")
        size = args.get("size", "1024x1024")
        style = args.get("style")  # e.g., "photorealistic", "watercolor", etc.
        if not prompt:
            return JSONResponse({"success": False, "error": "prompt required"}, status_code=400)
        try:
            # Call OpenAI Images API
            resp = requests.post(
                "https://api.openai.com/v1/images/generations",
                headers={
                    "Authorization": f"Bearer {OPENAI_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": "gpt-image-1",
                    # Include style by folding into prompt; supported params vary.
                    "prompt": f"{prompt}\n\nStyle: {style}" if style else prompt,
                    "size": size,
                },
                timeout=60,
            )
            log.info("/tool image.generate: status=%s", resp.status_code)
            if resp.status_code != 200:
                return JSONResponse({"success": False, "error": f"OpenAI error {resp.status_code}: {resp.text}"}, status_code=resp.status_code)
            data = resp.json()
            b64 = data.get("data", [{}])[0].get("b64_json")
            if not b64:
                return JSONResponse({"success": False, "error": "No image data returned"}, status_code=500)
            # Return as data URL to make client rendering trivial
            return JSONResponse({
                "success": True,
                "type": "image",
                "mime": "image/png",
                "data_url": f"data:image/png;base64,{b64}",
            })
        except Exception as e:
            return JSONResponse({"success": False, "error": str(e)}, status_code=500)

    if name == "file.save":
        # Save textual or tabular content to various formats
        try:
            base_dir = os.path.join(os.path.dirname(__file__), "..", "scratchpad", "saved")
            base_dir = os.path.abspath(base_dir)
            os.makedirs(base_dir, exist_ok=True)

            filename = (args.get("filename") or args.get("name") or "").strip() if args else ""
            fmt = (args.get("format") or args.get("type") or "").strip().lower() if args else ""
            content = args.get("content") if args else None
            rows = args.get("rows") if args else None  # list[dict]

            # Infer format from filename extension if needed
            if not fmt and filename:
                _, ext = os.path.splitext(filename)
                fmt = ext.lstrip(".").lower()

            # Normalize common aliases
            alias = {
                "doc": "docx",
                "ppt": "pptx",
                "xls": "xlsx",
                "markdown": "md",
                "text": "txt",
            }
            fmt = alias.get(fmt, fmt)

            if not fmt:
                return JSONResponse({"success": False, "error": "format required"}, status_code=400)

            # Default filename with date-stamp
            if not filename:
                stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
                filename = f"alita_output-{stamp}.{fmt or 'txt'}"
            # Ensure extension matches fmt
            root, ext = os.path.splitext(filename)
            if ext.lstrip(".").lower() != fmt:
                filename = f"{root}.{fmt}"

            # Safety: sanitize filename
            safe = re.sub(r"[^A-Za-z0-9._\- ]+", "_", filename).strip()
            file_path = os.path.join(base_dir, safe)

            # Handlers per format
            if fmt in ("txt", "md", "json", "py"):
                data = content
                if data is None and rows is not None:
                    import json as _json
                    data = _json.dumps(rows, indent=2)
                if data is None:
                    data = ""
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(str(data))
                return JSONResponse({"success": True, "type": "file", "path": file_path, "filename": safe, "format": fmt})

            if fmt == "pdf":
                from reportlab.lib.pagesizes import letter
                from reportlab.pdfgen import canvas
                text = str(content or "")
                c = canvas.Canvas(file_path, pagesize=letter)
                width, height = letter
                # simple text wrap
                y = height - 72
                for line in text.splitlines() or [""]:
                    c.drawString(72, y, line[:95])
                    y -= 14
                    if y < 72:
                        c.showPage()
                        y = height - 72
                c.save()
                return JSONResponse({"success": True, "type": "file", "path": file_path, "filename": safe, "format": fmt})

            if fmt == "docx":
                from docx import Document
                doc = Document()
                text = str(content or "")
                for para in (text.split("\n\n") if text else [""]):
                    doc.add_paragraph(para)
                doc.save(file_path)
                return JSONResponse({"success": True, "type": "file", "path": file_path, "filename": safe, "format": fmt})

            if fmt == "csv":
                import csv
                # rows expected as list of dicts; fallback: write content as single column
                if isinstance(rows, list) and rows:
                    fieldnames = sorted({k for r in rows if isinstance(r, dict) for k in r.keys()})
                    with open(file_path, "w", newline="", encoding="utf-8") as f:
                        w = csv.DictWriter(f, fieldnames=fieldnames)
                        w.writeheader()
                        for r in rows:
                            if isinstance(r, dict):
                                w.writerow({k: r.get(k) for k in fieldnames})
                else:
                    with open(file_path, "w", newline="", encoding="utf-8") as f:
                        w = csv.writer(f)
                        for line in str(content or "").splitlines():
                            w.writerow([line])
                return JSONResponse({"success": True, "type": "file", "path": file_path, "filename": safe, "format": fmt})

            if fmt == "xlsx":
                import pandas as pd
                if isinstance(rows, list) and rows:
                    df = pd.DataFrame(rows)
                else:
                    # store as single-column sheet
                    df = pd.DataFrame({"content": str(content or "").splitlines()})
                df.to_excel(file_path, index=False)
                return JSONResponse({"success": True, "type": "file", "path": file_path, "filename": safe, "format": fmt})

            if fmt == "pptx":
                from pptx import Presentation
                from pptx.util import Inches, Pt
                prs = Presentation()
                # Title + content slide
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = args.get("title") or "Alita Export"
                body = slide.placeholders[1]
                text = str(content or "")
                tf = body.text_frame
                tf.clear()
                first = True
                for line in text.splitlines() or [""]:
                    if first:
                        tf.text = line
                        first = False
                    else:
                        p = tf.add_paragraph(); p.text = line; p.level = 0
                prs.save(file_path)
                return JSONResponse({"success": True, "type": "file", "path": file_path, "filename": safe, "format": fmt})

            return JSONResponse({"success": False, "error": f"unsupported format: {fmt}"}, status_code=400)
        except Exception as e:
            log.exception("/tool file.save: exception: %s", e)
            return JSONResponse({"success": False, "error": str(e)}, status_code=500)

    if name == "search.web":
        # query (required), provider optional, max_results optional
        query = (args.get("query") or args.get("q") or "").strip()
        provider = (args.get("provider") or "").lower().strip() or ("tavily" if TAVILY_API_KEY else "duckduckgo")
        max_results = int(args.get("max_results", 6))
        if not query:
            return JSONResponse({"success": False, "error": "query required"}, status_code=400)
        try:
            # Try direct time resolver first
            results: list[dict[str, Any]] = []
            place = _extract_place_from_query(query)
            if place:
                mapped = _map_place_to_tz(place)
                if mapped:
                    place_label, tz = mapped
                    try:
                        tz_resp = requests.get(f"https://worldtimeapi.org/api/timezone/{tz}", timeout=15)
                        log.info("/tool search.web: time resolver tz=%s status=%s", tz, getattr(tz_resp, 'status_code', 'n/a'))
                        if tz_resp.status_code == 200:
                            tdata = tz_resp.json()
                            return JSONResponse({
                                "success": True,
                                "type": "time",
                                "place": place_label,
                                "tz": tz,
                                "datetime_iso": tdata.get("datetime"),
                                "utc_offset": tdata.get("utc_offset"),
                                "abbreviation": tdata.get("abbreviation"),
                                "source_url": f"https://worldtimeapi.org/api/timezone/{tz}",
                            })
                        # else: fall through to provider search below
                    except Exception:
                        pass

            # Provider web search
            if not place:
                log.info("/tool search.web: no time place extracted from query='%s'", query)
            if provider == "tavily":
                if not TAVILY_API_KEY:
                    provider = "duckduckgo"
                else:
                    resp = requests.post(
                        "https://api.tavily.com/search",
                        json={
                            "api_key": TAVILY_API_KEY,
                            "query": query,
                            "max_results": max_results,
                            "include_answer": False,
                        },
                        timeout=30,
                    )
                    log.info("/tool search.web: provider=tavily status=%s", resp.status_code)
                    if resp.status_code != 200:
                        return JSONResponse({"success": False, "error": f"Search error {resp.status_code}: {resp.text}"}, status_code=resp.status_code)
                    data = resp.json()
                    for item in data.get("results", [])[:max_results]:
                        results.append({
                            "title": item.get("title") or item.get("url"),
                            "url": item.get("url"),
                            "snippet": item.get("content") or item.get("snippet") or "",
                            "score": item.get("score"),
                            "source": "tavily",
                        })

            if provider == "duckduckgo":
                # DuckDuckGo Instant Answer API (no key). Not full SERP; gives abstracts and related topics.
                # https://api.duckduckgo.com/?q=...&format=json
                def ddg_request():
                    return requests.get(
                        "https://api.duckduckgo.com/",
                        params={"q": query, "format": "json", "no_html": 1, "skip_disambig": 1},
                        timeout=20,
                    )
                resp = ddg_request()
                log.info("/tool search.web: provider=duckduckgo status=%s", resp.status_code)
                if resp.status_code != 200:
                    # Treat 202 (offline/test backends) as empty-success rather than loud error
                    if resp.status_code == 202:
                        data = {}
                    else:
                        # Retry once for transient network issues
                        try:
                            resp2 = ddg_request()
                        except Exception:
                            resp2 = None
                        if resp2 is not None and resp2.status_code == 200:
                            data = resp2.json()
                        else:
                            body = resp.text if isinstance(resp.text, str) else str(resp.text)
                            body = (body[:500] + "…") if len(body) > 500 else body
                            log.warning("/tool search.web: duckduckgo error status=%s body=%s", resp.status_code, body)
                            return JSONResponse(
                                {"success": False, "error": f"provider=duckduckgo status={resp.status_code} body={body}"},
                                status_code=resp.status_code,
                            )
                else:
                    data = resp.json()
                # Primary abstract
                abstract_text = (data.get("AbstractText") or "").strip()
                abstract_url = (data.get("AbstractURL") or "").strip() or (data.get("AbstractSource") or "")
                if abstract_text and abstract_url:
                    results.append({
                        "title": data.get("Heading") or abstract_url,
                        "url": abstract_url,
                        "snippet": abstract_text,
                        "source": "duckduckgo",
                    })
                # Related topics (flatten one level)
                def iter_related(items):
                    for it in items or []:
                        if "FirstURL" in it:
                            yield it
                        for sub in it.get("Topics", []) or []:
                            if "FirstURL" in sub:
                                yield sub
                for it in iter_related(data.get("RelatedTopics")):
                    results.append({
                        "title": it.get("Text") or it.get("FirstURL"),
                        "url": it.get("FirstURL"),
                        "snippet": it.get("Text") or "",
                        "source": "duckduckgo",
                    })
                results = results[:max_results]

            log.info("/tool search.web: results=%d provider=%s", len(results), provider)
            # If DuckDuckGo returned an empty 202 and this looks like a time query, try time resolver again as a fallback
            if not results:
                place_fb = _extract_place_from_query(query)
                if place_fb:
                    mapped = _map_place_to_tz(place_fb)
                    if mapped:
                        place_label, tz = mapped
                        try:
                            tz_resp = requests.get(f"https://worldtimeapi.org/api/timezone/{tz}", timeout=15)
                            log.info("/tool search.web: fallback time resolver tz=%s status=%s", tz, getattr(tz_resp, 'status_code', 'n/a'))
                            if tz_resp.status_code == 200:
                                tdata = tz_resp.json()
                                return JSONResponse({
                                    "success": True,
                                    "type": "time",
                                    "place": place_label,
                                    "tz": tz,
                                    "datetime_iso": tdata.get("datetime"),
                                    "utc_offset": tdata.get("utc_offset"),
                                    "abbreviation": tdata.get("abbreviation"),
                                    "source_url": f"https://worldtimeapi.org/api/timezone/{tz}",
                                })
                        except Exception:
                            pass

            return JSONResponse({
                "success": True,
                "type": "search",
                "query": query,
                "results": results,
                "provider": provider,
            })
        except Exception as e:
            log.exception("/tool search.web: exception: %s", e)
            return JSONResponse({"success": False, "error": str(e)}, status_code=500)

    return JSONResponse({"success": False, "error": f"unknown tool: {name}"}, status_code=400)

if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("ALITA_REALTIME_PORT", "8787"))
    uvicorn.run("app.realtime_proxy:app", host="127.0.0.1", port=port, reload=False)
